# Copyright 2026 Google LLC
#
# Licensed under the Apache License the License is distributed on an "AS IS" BASIS,# Licensed under the Apache License, Version 2.0 (the "License");
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import asyncio
import json
import os
import re
import tempfile
import uuid
from typing import Any, Optional

from google.adk.tools.tool_context import ToolContext
from google.genai import types
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT

from ..shared_libraries.config import (
    DEFAULT_TEMPLATE_URI,
    get_logger,
)
from ..shared_libraries.models import CoverSpec, DeckSpec, SlideSpec
from ..shared_libraries.utils import _insert_image
from .artifact_utils import get_gcs_file_as_local_path, save_presentation
from .pptx_editor import _insert_visual_into_slide
from .visual_generator import generate_visual


def _sanitize_filename(name: str) -> str:
    """Make a filesystem + GCS-friendly filename."""
    name = name.strip()
    name = re.sub(r"[^\w\- ]+", "", name)  # remove illegal chars
    name = re.sub(r"\s+", "_", name)
    return name[:80] if name else "generated"


def get_smart_layout(prs: Presentation, requested_name: str):
    """
    Maps a requested layout name to the best available layout in the current presentation.
    """
    log = get_logger("layout_mapper")
    requested_name = (requested_name or "").lower()

    # Avoid "Title and Chart" layouts because they often require embedded chart data
    if "chart" in requested_name:
        log.info(
            f"Overriding layout request '{requested_name}' to 'Title and Image' for better compatibility."
        )
        requested_name = "title and image"

    layouts = prs.slide_layouts

    # 1) Exact match
    for layout in layouts:
        if layout.name.lower() == requested_name:
            return layout

    # 2) Keyword mapping
    mapping = [
        ("two", ["two content", "comparison", "side by side", "dual", "split"]),
        (
            "image",
            ["image", "picture", "visual", "photo", "graphic", "title and image", "image grid"],
        ),
        ("quote", ["quote", "testimonial", "statement"]),
        ("agenda", ["agenda", "toc", "roadmap", "contents"]),
        ("section", ["section header", "divider", "transition"]),
        ("content", ["title and content", "content slide", "standard body", "bullet", "subhead", "left"]),
        ("closing", ["closing", "thank you", "end", "contact"]),
        ("title", ["title slide", "cover", "intro", "opening", "only"]),
    ]

    for concept, keywords in mapping:
        if concept in requested_name:
            for layout in layouts:
                if any(k in layout.name.lower() for k in keywords):
                    log.info(f"Mapped '{requested_name}' -> '{layout.name}'")
                    return layout

    # 3) Fallbacks
    try:
        if "title" in requested_name or "cover" in requested_name:
            for layout in layouts:
                if "title" in layout.name.lower() and "content" not in layout.name.lower():
                    return layout
            return layouts[0]

        return layouts[1] if len(layouts) > 1 else layouts[0]
    except Exception:
        return layouts[0]


def _find_title_placeholder(slide):
    # title placeholder or shapes.title
    try:
        if slide.shapes.title:
            return slide.shapes.title
    except Exception:
        pass
    for shape in getattr(slide, "placeholders", []):
        try:
            if shape.placeholder_format.type in (PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE):
                return shape
        except Exception:
            continue
    return None


def _find_subtitle_placeholder(slide):
    for shape in getattr(slide, "placeholders", []):
        try:
            if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                return shape
        except Exception:
            continue
    return None


def _find_content_placeholders(slide):
    phs = []
    for shape in getattr(slide, "placeholders", []):
        try:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER_TYPE.BODY,
                PP_PLACEHOLDER_TYPE.OBJECT,
                PP_PLACEHOLDER_TYPE.PICTURE,
            ):
                phs.append(shape)
        except Exception:
            continue
    # sort left-to-right
    try:
        phs = sorted(phs, key=lambda p: p.left)
    except Exception:
        pass
    return phs


def _find_first_table(slide):
    for shape in slide.shapes:
        try:
            if shape.has_table:
                return shape.table
        except Exception:
            continue
    return None


def _apply_bullets(text_frame, bullets):
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    if not bullets:
        return

    for i, bullet_text in enumerate(bullets):
        is_sub = (
            bullet_text.startswith("  ")
            or bullet_text.startswith("\t")
            or bullet_text.startswith("- ")
        )
        clean_text = bullet_text.strip(" \t-•*")

        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.level = 1 if is_sub else 0

        # treat **bold** markers as emphasis, but keep simple
        parts = clean_text.split("**")
        for j, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            r.text = part
            if j % 2 != 0:
                r.font.bold = True


def _apply_table(table, table_dict: dict):
    """
    table_dict supports:
      {"headers": [...], "rows": [[...],[...]]}
      OR {"data": [[...],[...]]}
    """
    if not table or not table_dict:
        return

    data = None
    headers = table_dict.get("headers")
    rows = table_dict.get("rows")

    if headers is not None:
        data = [headers] + (rows or [])
    elif "data" in table_dict:
        data = table_dict.get("data")

    if not data:
        return

    max_rows = len(table.rows)
    max_cols = len(table.columns)

    for r_i, row in enumerate(data):
        if r_i >= max_rows:
            break
        for c_i, val in enumerate(row):
            if c_i >= max_cols:
                break
            table.cell(r_i, c_i).text = "" if val is None else str(val)


def _render_slide_content(prs: Presentation, slide, spec_obj, is_cover: bool = False):
    """
    Populates a slide based on:
      - title/subtitle
      - bullets
      - table (if available)
      - image_data/image_file_path
    """
    # 1) Title & subtitle
    title_ph = _find_title_placeholder(slide)
    subtitle_ph = _find_subtitle_placeholder(slide)

    if title_ph and getattr(spec_obj, "title", None):
        title_ph.text = str(spec_obj.title).replace("**", "")
        # center title on cover
        if is_cover:
            try:
                if title_ph.text_frame.paragraphs:
                    title_ph.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            except Exception:
                pass

    if subtitle_ph and getattr(spec_obj, "subhead", None):
        subtitle_ph.text = str(spec_obj.subhead).replace("**", "")

    if is_cover:
        return

    # 2) If table is present in slide and spec has table data -> fill table
    table = _find_first_table(slide)
    table_dict = None
    if hasattr(spec_obj, "table") and getattr(spec_obj, "table"):
        table_dict = getattr(spec_obj, "table")
    elif isinstance(spec_obj, dict) and spec_obj.get("table"):
        table_dict = spec_obj.get("table")

    if table and isinstance(table_dict, dict):
        _apply_table(table, table_dict)

    # 3) Text & visuals
    bullets = getattr(spec_obj, "bullets", None)
    has_bullets = bool(bullets)
    image_source = getattr(spec_obj, "image_data", None) or getattr(spec_obj, "image_file_path", None)

    content_phs = _find_content_placeholders(slide)

    # both text + image
    if has_bullets and image_source:
        if len(content_phs) >= 2:
            _apply_bullets(content_phs[0].text_frame, bullets)
            _insert_visual_into_slide(prs, slide, image_source, target_placeholder=content_phs[1])
        elif len(content_phs) == 1:
            _apply_bullets(content_phs[0].text_frame, bullets)
            _insert_visual_into_slide(prs, slide, image_source, target_placeholder=None)

    # only text
    elif has_bullets and not image_source:
        if len(content_phs) >= 1:
            _apply_bullets(content_phs[0].text_frame, bullets)

    # only image
    elif (not has_bullets) and image_source:
        _insert_visual_into_slide(prs, slide, image_source, target_placeholder=content_phs[0] if content_phs else None)


async def render_deck_from_spec(
    spec_dict: dict,
    tool_context: ToolContext,
    template_pptx: str,
) -> str:
    """
    Renders a presentation from a spec using a REQUIRED template_pptx (local path).
    IMPORTANT DQA CHANGE:
      - We DO NOT delete the template slides.
      - We populate existing template slides by index.
      - If spec has more slides than template, we append new slides using best layouts.
    Returns: local generated pptx path or "Error: ..."
    """
    log = get_logger("render_deck_from_spec")

    try:
        if not template_pptx or not os.path.exists(template_pptx):
            return "Error: No valid template provided."

        prs = Presentation(template_pptx)

        # Normalize spec
        if isinstance(spec_dict.get("slides"), dict):
            spec_dict["slides"] = list(spec_dict["slides"].values())
        if "closing_title" not in spec_dict:
            spec_dict["closing_title"] = "Thank You"

        # Validate using existing model
        validated_spec = DeckSpec(**spec_dict)

        # Ensure we have at least one slide in the template
        if len(prs.slides) == 0:
            # If template is empty, create a cover slide
            cover_slide = prs.slides.add_slide(get_smart_layout(prs, "Title Slide"))
        else:
            cover_slide = prs.slides[0]

        # Render cover onto first slide
        _render_slide_content(prs, cover_slide, validated_spec.cover, is_cover=True)

        # Render body slides into existing template slides starting from slide 2
        template_slide_count = len(prs.slides)
        body_specs = validated_spec.slides or []

        for i, s_spec in enumerate(body_specs):
            target_index = i + 1  # because slide 0 is cover
            if target_index < template_slide_count:
                slide = prs.slides[target_index]
            else:
                # append new slide if spec has more than template
                slide = prs.slides.add_slide(get_smart_layout(prs, s_spec.layout_name))

            _render_slide_content(prs, slide, s_spec, is_cover=False)

            # speaker notes (no citations in DQA flow by default)
            if getattr(s_spec, "speaker_notes", None):
                try:
                    slide.notes_slide.notes_text_frame.text = s_spec.speaker_notes
                except Exception:
                    pass

        # Optional: closing slide if template has one extra slot at end and spec provides closing title
        # If template already has an extra slide, set its title to closing_title.
        # Otherwise add a closing slide.
        closing_title = validated_spec.closing_title or "Thank You"
        if len(prs.slides) >= (len(body_specs) + 2):
            # there is at least one spare slide after expected content; treat last slide as closing
            closing_slide = prs.slides[-1]
        else:
            closing_slide = prs.slides.add_slide(get_smart_layout(prs, "Closing Slide"))

        # set closing title if there is a title shape
        t_ph = _find_title_placeholder(closing_slide)
        if t_ph:
            try:
                t_ph.text = closing_title
            except Exception:
                pass

        # Save to temp output path
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            prs.save(tmp.name)
            return tmp.name

    except Exception as e:
        log.error(f"Render failed: {e}", exc_info=True)
        return f"Error: Render failed. {e}"


async def generate_and_render_deck(
    tool_context: ToolContext,
    deck_spec: dict | None = None,
    spec_artifact_name: str | None = None,
    template_path: str | None = None,
    template_gcs_uri: str | None = None,
) -> dict:
    """
    Orchestrates deck rendering from an existing deck_spec.
    DQA changes:
      - NO dependency on research.
      - Template should come from TEMPLATE_BUCKET selection (passed in template_gcs_uri) or local template_path.
      - Output saving uses save_presentation() which uploads to OUTPUT_BUCKET by default.
    """
    log = get_logger("generate_and_render_deck_tool")

    try:
        spec_dict = deck_spec

        # 1) Load spec from session state if not provided
        if not spec_dict and not spec_artifact_name:
            spec_dict = tool_context.state.get("current_deck_spec")
            if spec_dict:
                log.info("Loaded DeckSpec from session state.")

        # 2) Fallback: load spec from named artifact
        if not spec_dict and spec_artifact_name:
            log.info(f"Loading DeckSpec from artifact: '{spec_artifact_name}'")
            artifact = await tool_context.load_artifact(spec_artifact_name)

            if not artifact:
                log.warning(f"Artifact '{spec_artifact_name}' not found. Waiting 2s...")
                await asyncio.sleep(2.0)
                artifact = await tool_context.load_artifact(spec_artifact_name)

            if artifact:
                spec_json = artifact.inline_data.data if isinstance(artifact, types.Part) else artifact
                if isinstance(spec_json, (bytes, bytearray)):
                    spec_dict = json.loads(spec_json.decode("utf-8"))
                elif isinstance(spec_json, str):
                    spec_dict = json.loads(spec_json)
                else:
                    spec_dict = spec_json

        if not spec_dict:
            return {
                "status": "Failed",
                "message": "No active presentation plan found. Provide deck_spec or generate an outline first.",
            }

        # 3) Resolve template path (template_path or template_gcs_uri or DEFAULT_TEMPLATE_URI)
        working_template = template_path
        if working_template and not os.path.exists(working_template):
            working_template = None

        if not working_template and template_gcs_uri:
            log.info("Downloading selected template from GCS...")
            working_template = await get_gcs_file_as_local_path(template_gcs_uri)

        if not working_template:
            # fallback (kept for compatibility) - but in DQA flow you should pass selected template.
            if DEFAULT_TEMPLATE_URI:
                log.warning("No template provided; falling back to DEFAULT_TEMPLATE_URI.")
                working_template = await get_gcs_file_as_local_path(DEFAULT_TEMPLATE_URI)

        if not working_template or str(working_template).startswith("Error"):
            return {"status": "Failed", "message": f"Template could not be resolved: {working_template}"}

        # 4) Validate spec dict
        if isinstance(spec_dict.get("slides"), dict):
            spec_dict["slides"] = list(spec_dict["slides"].values())
        if "closing_title" not in spec_dict:
            spec_dict["closing_title"] = "Thank You"

        validated_spec = DeckSpec(**spec_dict)

        # 5) Optional visuals (ONLY if visual_prompt exists; no web search)
        # Hard limit to prevent overload
        hard_limit = 5
        visuals_kept = 0
        for slide in validated_spec.slides:
            if getattr(slide, "visual_prompt", None):
                if visuals_kept < hard_limit:
                    visuals_kept += 1
                    # keep layout compatible for image usage
                    if slide.layout_name not in ["Title and Image", "Two Content", "Comparison"]:
                        slide.layout_name = "Title and Image"
                else:
                    slide.visual_prompt = None

        tasks = []
        slides_with_visuals = []
        for item in [validated_spec.cover] + validated_spec.slides:
            if getattr(item, "visual_prompt", None):
                tasks.append(asyncio.create_task(asyncio.wait_for(generate_visual(item.visual_prompt), timeout=60.0)))
                slides_with_visuals.append(item)

        if tasks:
            images = await asyncio.gather(*tasks, return_exceptions=True)
            for s, img in zip(slides_with_visuals, images):
                if not isinstance(img, Exception):
                    s.image_data = img

        # 6) Render to local file
        local_path = await render_deck_from_spec(
            validated_spec.model_dump(),
            tool_context,
            working_template,
        )

        if str(local_path).startswith("Error:"):
            return {"status": "Failed", "message": local_path}

        # 7) Save output (uploads to OUTPUT_BUCKET by default in our updated artifact_utils)
        base = _sanitize_filename(validated_spec.cover.title or "proposal")
        out_name = f"{base}_{uuid.uuid4().hex[:6]}.pptx"

        msg = await save_presentation(tool_context, out_name, local_path)
        try:
            os.remove(local_path)
        except Exception:
            pass

        return {"status": "Success", "message": msg}

    except Exception as e:
        log.error(f"Generation failed: {e}", exc_info=True)
        return {"status": "Failed", "message": str(e)}

# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
