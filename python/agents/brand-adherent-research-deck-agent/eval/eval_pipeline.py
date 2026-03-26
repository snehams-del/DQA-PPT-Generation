import asyncio
import json
import os
import re
import tempfile
from io import BytesIO

from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation

# Load env before importing app to ensure we can manipulate it if needed
load_dotenv()

# Force InMemory services for evaluation by unsetting the bucket name
os.environ["GCP_STAGING_BUCKET"] = ""

from presentation_agent.agent import PresentationExpertApp

def create_valid_mock_pptx():
    """Generates bytes for a valid 5-slide PowerPoint file."""
    prs = Presentation()
    for i in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
        slide.shapes.title.text = f"Slide {i} Title"
        slide.placeholders[1].text = f"This is the body content for slide {i}. It contains useful data for testing."
    
    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()

async def mock_upload_artifact(app, session_id, filename):
    """Mocks an 'upload' by saving a valid PPTX into the agent's session store."""
    content = create_valid_mock_pptx()
        
    artifact = types.Part(
        inline_data=types.Blob(
            data=content,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    )
    
    # Save the artifact using the EXACT keyword arguments from ADK signature
    await app._runner.artifact_service.save_artifact(
        app_name="presentation_expert",
        user_id="evaluator",
        filename=filename, 
        artifact=artifact,
        session_id=session_id
    )
    return filename


async def evaluate_scenario(client, app, test_case, session_id):
    """Runs a single test scenario and calculates metrics."""
    print(f"\n--- Running Scenario: {test_case['name']} ---")
    
    # Pre-load artifacts if the test case requires them
    if "artifacts" in test_case:
        for filename in test_case["artifacts"]:
            await mock_upload_artifact(app, session_id, filename)
            print(f"  [Setup] Mocked upload of valid PPTX: {filename}")

    actual_tool_calls = []
    final_response = ""
    generated_deck_spec = None

    # --- TURN 1: Initial Request ---
    print(f"Turn 1: Requesting... '{test_case['prompt']}'")
    message = types.Content(
        role="user", parts=[types.Part.from_text(text=test_case["prompt"])]
    )
    
    # Use the real application's runner
    events = app._runner.run(
        user_id="evaluator", session_id=session_id, new_message=message
    )
    for event in events:
        if event.content and event.content.parts:
            for part in event.content.parts:
                if part.function_call:
                    tool_name = part.function_call.name
                    actual_tool_calls.append(tool_name)
                    print(f"  [Tool Call]: {tool_name}")
                if part.text:
                    final_response += part.text

    # --- TURN 2: Follow-up / Approval (if needed) ---
    if test_case.get("requires_approval", False):
        print("Turn 2: Sending 'Approve' to trigger rendering...")
        approve_msg = types.Content(
            role="user", parts=[types.Part.from_text(text="Approve and render now.")]
        )
        events2 = app._runner.run(
            user_id="evaluator", session_id=session_id, new_message=approve_msg
        )
        for event in events2:
            if event.content and event.content.parts:
                for part in event.content.parts:
                    if part.text:
                        final_response += part.text
                    if part.function_call:
                        tool_name = part.function_call.name
                        actual_tool_calls.append(tool_name)
                        print(f"  [Tool Call]: {tool_name}")
                        # Robustly capture deck_spec from either tool
                        if tool_name in ["save_deck_spec", "generate_and_render_deck"]:
                            args = part.function_call.args
                            if "deck_spec" in args:
                                generated_deck_spec = args["deck_spec"]

    # --- Metrics Calculation ---

    # A. Tool Trajectory Score
    expected = test_case["expected_tools"]
    
    # Check for strict sequential ordering if required (e.g., for editing workflows)
    if test_case.get("strict_order", False):
        filtered_actual = [call for call in actual_tool_calls if call in expected]
        deduped_actual = []
        for call in filtered_actual:
            if not deduped_actual or deduped_actual[-1] != call:
                deduped_actual.append(call)
        
        trajectory_match = False
        if len(deduped_actual) >= len(expected):
             for i in range(len(deduped_actual) - len(expected) + 1):
                 if deduped_actual[i:i+len(expected)] == expected:
                     trajectory_match = True
                     break
        tool_trajectory_avg_score = 1.0 if trajectory_match else 0.0
        print(f"  [Metric] Strict Order Match: {trajectory_match}")
    else:
        intersection = set(expected).intersection(set(actual_tool_calls))
        tool_trajectory_avg_score = len(intersection) / max(len(expected), 1)

    # B. Response Match Score
    judge_match_prompt = f"""
    Evaluate how well the agent's final response matches the intent.
    User Intent: "{test_case['prompt']}"
    Target Keywords: {test_case['target_keywords']}
    Agent Response: "{final_response}"
    Score from 0.0 to 1.0. Return ONLY the float number.
    """
    try:
        match_resp = client.models.generate_content(
            model="gemini-2.5-flash", contents=judge_match_prompt
        )
        response_match_score = float(
            re.search(r"(\d+\.\d+|\d+)", match_resp.text).group(1)
        )
    except Exception:
        response_match_score = 0.0

    # C. Constraint Compliance Score
    constraint_compliance_score = 0.0
    score = 0
    total_constraints = 0

    if "constraints" in test_case:
        constraints = test_case["constraints"]
        
        if "slide_count" in constraints and generated_deck_spec:
            total_constraints += 1
            if len(generated_deck_spec.get("slides", [])) == constraints["slide_count"]:
                score += 1
                
        if "required_terms" in constraints and generated_deck_spec:
            total_constraints += 1
            deck_text = json.dumps(generated_deck_spec).lower()
            if all(term.lower() in deck_text for term in constraints["required_terms"]):
                score += 1

        if "has_citations" in constraints:
            total_constraints += 1
            # Look for explicit source tags, raw URLs, OR standard markdown links
            has_explicit_source = bool(re.search(r"\(Source: <https?://.*?>\)", final_response))
            has_raw_url = bool(re.search(r"https?://[^\s]+", final_response))
            has_markdown_url = bool(re.search(r"\[.*?\]\(https?://.*?\)", final_response))
            
            if has_explicit_source or has_raw_url or has_markdown_url:
                score += 1
            else:
                print(f"  [Debug] Citation check failed. final_response: {final_response[:200]}...")

    constraint_compliance_score = (score / total_constraints) if total_constraints > 0 else 1.0

    results = {
        "scenario": test_case["name"],
        "tool_trajectory_avg_score": round(tool_trajectory_avg_score, 2),
        "response_match_score": round(response_match_score, 2),
        "constraint_compliance_score": round(constraint_compliance_score, 2),
    }
    return results


async def run_evaluation_with_metrics():
    """
    Standard Evaluation Pipeline testing Creation, Editing, and Research capabilities.
    """
    print("=========================================")
    print("🚀 Starting Multi-Scenario Eval Pipeline 🚀")
    print("=========================================\n")

    project_id = os.environ.get("GOOGLE_CLOUD_PROJECT")
    location = os.environ.get("GOOGLE_CLOUD_LOCATION", "us-central1")

    if not project_id:
        print("Skipping: GOOGLE_CLOUD_PROJECT not set.")
        return

    client = genai.Client(vertexai=True, project=project_id, location=location)

    # Initialize the Real App
    app = PresentationExpertApp()

    # 1. Define Test Scenarios
    test_scenarios = [
        {
            "name": "Standard Creation Workflow",
            "prompt": "Create a 2-slide presentation about why Cloud Computing is good. You must mention 'AWS' and 'Azure'. No external tools.",
            "requires_approval": True,
            "expected_tools": [
                "get_gcs_file_as_local_path",
                "inspect_template_style",
                "generate_and_save_outline",
                "batch_generate_slides",
                "generate_and_render_deck",
            ],
            "target_keywords": ["presentation", "ready"],
            "constraints": {"slide_count": 2, "required_terms": ["AWS", "Azure"]},
        },
        {
            "name": "Editing: Index Shifting",
            "prompt": "I have uploaded 'draft.pptx'. First, delete slide 2. Then, add a new slide at the end summarizing our Q3 goals.",
            "artifacts": ["draft.pptx"],
            "requires_approval": False,
            "strict_order": False,
            "expected_tools": [
                "get_artifact_as_local_path",
                "delete_slide",
                "read_presentation_outline",
                "add_slide_to_end",
            ],
            "target_keywords": ["deleted", "added"],
            "constraints": {},
        },
        {
            "name": "Editing: Modify Text",
            "prompt": "Using 'draft.pptx', summarize the text on slide 3 so it is punchier. Keep the same title.",
            "artifacts": ["draft.pptx"],
            "requires_approval": False,
            "expected_tools": [
                "get_artifact_as_local_path",
                "extract_slide_content",
                "edit_slide_text",
                "read_presentation_outline",
            ],
            "target_keywords": ["updated", "slide 3"],
            "constraints": {},
        },
        {
            "name": "Editing: Visuals & Layout",
            "prompt": "In 'draft.pptx', generate a new chart showing 50% growth and replace the visual on slide 4. Once replaced, move that element slightly to the right.",
            "artifacts": ["draft.pptx"],
            "requires_approval": False,
            "strict_order": False,
            "expected_tools": [
                "get_artifact_as_local_path",
                "generate_visual",
                "replace_slide_visual",
                "update_element_layout",
                "read_presentation_outline",
            ],
            "target_keywords": ["replaced", "moved"],
            "constraints": {},
        },
        {
            "name": "Editing: Cross-Deck Extraction & Save",
            "prompt": "Extract the text from slide 5 of 'reference.pptx' and add it as a new slide to my current working file 'draft.pptx'. That's all, save the final presentation as 'final_deck.pptx'.",
            "artifacts": ["draft.pptx", "reference.pptx"],
            "requires_approval": False,
            "expected_tools": [
                "get_artifact_as_local_path",
                "extract_slide_content",
                "add_slide_to_end",
                "save_presentation",
            ],
            "target_keywords": ["final_deck.pptx", "saved"],
            "constraints": {},
        },
        {
            "name": "Deep Research & Citation Compliance",
            "prompt": "Do a deep research analysis on the 10-year tech sector growth strategy for Irvine. Provide a detailed markdown summary with citations.",
            "requires_approval": False,
            "expected_tools": [
                "research_specialist"
            ],
            "target_keywords": ["Irvine", "growth"],
            "constraints": {
                "has_citations": True
            },
        }
    ]

    # 2. Run Evaluations
    all_results = []

    for i, test_case in enumerate(test_scenarios):
        session_id = f"eval_pipeline_session_{i}"
        await app._runner.session_service.create_session(
            app_name="presentation_expert", user_id="evaluator", session_id=session_id
        )
        
        try:
             res = await evaluate_scenario(client, app, test_case, session_id)
             all_results.append(res)
        except Exception as e:
             print(f"Error evaluating scenario '{test_case['name']}': {e}")

    # 3. Aggregate and Format Results
    print("\n" + "=" * 50)
    print("📊 FINAL EVALUATION METRICS BY SCENARIO:")
    print(json.dumps(all_results, indent=4))
    
    if all_results:
        avg_traj = sum(r["tool_trajectory_avg_score"] for r in all_results) / len(all_results)
        avg_match = sum(r["response_match_score"] for r in all_results) / len(all_results)
        avg_const = sum(r["constraint_compliance_score"] for r in all_results) / len(all_results)
        
        print("\n--- OVERALL AVERAGES ---")
        print(f"Tool Trajectory: {avg_traj:.2f}")
        print(f"Response Match:  {avg_match:.2f}")
        print(f"Constraint Comp: {avg_const:.2f}")
    
    print("=" * 50)

if __name__ == "__main__":
    asyncio.run(run_evaluation_with_metrics())
