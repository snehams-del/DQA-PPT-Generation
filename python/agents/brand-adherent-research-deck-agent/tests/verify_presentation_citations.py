# Copyright 2026 Google LLC
import asyncio
import os
from pptx import Presentation
from google.adk.tools.tool_context import ToolContext
from presentation_agent.tools.presentation_orchestrator import render_deck_from_spec
from presentation_agent.shared_libraries.models import DeckSpec, SlideSpec, CoverSpec

async def test_citations_in_speaker_notes_only():
    """
    Verifies that URLs are NOT in bullet points but ARE in speaker notes.
    """
    # 1. Setup a DeckSpec with citations
    deck_spec = {
        "cover": {"title": "Citation Test"},
        "slides": [
            {
                "title": "Slide 1",
                "bullets": ["Fact 1", "Fact 2"],
                "citations": ["https://example.com/source1"],
                "speaker_notes": "Main talk track."
            }
        ],
        "closing_title": "End"
    }
    
    # 2. Render to a temp file
    # We need a dummy template. Let's try to use the default one if it exists or mock it.
    template_path = "docs/Proposal_Template.pptx" 
    if not os.path.exists(template_path):
        print(f"Skipping test: Template not found at {template_path}")
        return

    out_pptx = "tests/citation_test_output.pptx"
    tool_context = None # Not strictly needed for the render logic we are testing
    
    result_path = await render_deck_from_spec(deck_spec, out_pptx, tool_context, template_pptx=template_path)
    
    assert not result_path.startswith("Error")
    
    # 3. Inspect the rendered PPTX
    prs = Presentation(result_path)
    slide = prs.slides[1] # Slide 1 (after cover)
    
    # Check bullets (Shape 1 usually)
    bullets_text = ""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            bullets_text += shape.text_frame.text
            
    assert "https://example.com/source1" not in bullets_text
    print("✅ Verified: URL is NOT in bullet points.")
    
    # Check speaker notes
    notes = slide.notes_slide.notes_text_frame.text
    assert "https://example.com/source1" in notes
    assert "Citations:" in notes
    print("✅ Verified: URL IS in speaker notes.")
    
    # Cleanup
    if os.path.exists(result_path):
        os.remove(result_path)

if __name__ == "__main__":
    asyncio.run(test_citations_in_speaker_notes_only())
