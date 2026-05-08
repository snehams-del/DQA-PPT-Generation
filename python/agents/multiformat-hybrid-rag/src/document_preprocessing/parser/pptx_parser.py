# Copyright 2026 Google LLC
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     https://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""PPTX → Markdown parser.

Converts to PDF via LibreOffice, extracts raw text per slide via python-pptx,
then delegates to parse_pdf_pages for parallel Gemini processing.
"""

from __future__ import annotations

import logging

from pptx import Presentation

from src.document_preprocessing.parser.convert import to_pdf
from src.document_preprocessing.parser.pdf_parser import (
    _split_pdf_pages,
    parse_pdf_pages,
)

logger = logging.getLogger(__name__)


def _extract_slide_texts(file_path: str) -> list[str]:
    """Extract raw text from each slide using python-pptx."""
    prs = Presentation(file_path)
    slide_texts = []
    for slide in prs.slides:
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        slide_texts.append("\n".join(lines))
    return slide_texts


def parse(file_path: str, genai_client=None, max_workers: int = 8) -> str:
    """Extract markdown from a .pptx file.

    1. Convert to PDF via LibreOffice
    2. Split PDF into individual pages
    3. Extract raw text per slide via python-pptx
    4. Send each page PDF + raw text to Gemini in parallel
    """
    if genai_client is None:
        raise ValueError("genai_client is required for PPTX parsing")

    pdf_path = to_pdf(file_path)
    pages, _ = _split_pdf_pages(pdf_path)
    texts = _extract_slide_texts(file_path)

    return parse_pdf_pages(pages, genai_client, texts=texts, max_workers=max_workers)


def quick_raw_text(file_path: str) -> str:
    """Cheap raw-text extraction (no LibreOffice, no LLM).

    python-pptx works directly on .pptx files. For .ppt (legacy), call
    via the ppt_parser path which converts to .pptx first.
    """
    return "\n\n".join(_extract_slide_texts(file_path))
